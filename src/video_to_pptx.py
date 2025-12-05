"""
------------------------------------------------------------------------------
 Application: Video to PPTX Converter
 Author: tanbaycu
 Copyright (c) 2025 tanbaycu. All rights reserved.
 
 Description: Converts video frames into a PowerPoint presentation based on visual changes.
------------------------------------------------------------------------------
"""
import cv2
import os
import numpy as np
from pptx import Presentation
from pptx.util import Inches
import argparse
import time

def mse(imageA, imageB):
    err = np.sum((imageA.astype("float") - imageB.astype("float")) ** 2)
    err /= float(imageA.shape[0] * imageA.shape[1])
    return err

def convert_mp4_to_pptx(video_path, output_pptx_path, threshold=2000, sample_rate_sec=1.0):
    print(f"Starting conversion for: {video_path}")
    
    if not os.path.exists(video_path):
        print(f"Error: File not found {video_path}")
        return

    cap = cv2.VideoCapture(video_path)
    if not cap.isOpened():
        print("Error opening video stream")
        return

    fps = cap.get(cv2.CAP_PROP_FPS)
    total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
    duration = total_frames / fps
    width = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
    height = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
    
    print(f"Video: {width}x{height}, FPS={fps}, Duration={duration:.2f}s")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    last_saved_frame = None
    slide_count = 0
    temp_dir = "temp_slides"
    if not os.path.exists(temp_dir):
        os.makedirs(temp_dir)
        
    frame_interval = int(fps * sample_rate_sec)
    current_frame_idx = 0
    
    while True:
        cap.set(cv2.CAP_PROP_POS_FRAMES, current_frame_idx)
        ret, frame = cap.read()
        if not ret:
            break
            
        gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
        small_gray = cv2.resize(gray, (50, 50))
        
        is_new_slide = False
        
        if last_saved_frame is None:
            is_new_slide = True
        else:
            similarity_error = mse(small_gray, last_saved_frame)
            if similarity_error > threshold:
                is_new_slide = True
                print(f"New slide at {current_frame_idx/fps:.1f}s (Diff: {similarity_error:.2f})")
        
        if is_new_slide:
            image_path = os.path.join(temp_dir, f"slide_{slide_count}.png")
            
            # Upscale to 4K logic (simplified 3x upscale)
            h, w = frame.shape[:2]
            high_res = cv2.resize(frame, (w * 3, h * 3), interpolation=cv2.INTER_CUBIC)
            enhanced = cv2.detailEnhance(high_res, sigma_s=10, sigma_r=0.15)
            
            cv2.imwrite(image_path, enhanced)
            
            blank_slide_layout = prs.slide_layouts[6] 
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
            
            last_saved_frame = small_gray
            slide_count += 1
            
        current_frame_idx += frame_interval
        if current_frame_idx > total_frames:
            break

    cap.release()
    
    print(f"Saving {slide_count} slides to {output_pptx_path}...")
    try:
        prs.save(output_pptx_path)
        print("Done!")
    except PermissionError:
        timestamp = int(time.time())
        alt_path = f"output_{timestamp}.pptx"
        prs.save(alt_path)
        print(f"File open error. Saved to {alt_path} instead.")
    
    # Cleanup
    for file in os.listdir(temp_dir):
        os.remove(os.path.join(temp_dir, file))
    os.rmdir(temp_dir)

def main():
    parser = argparse.ArgumentParser(description="Convert MP4 video to PPTX presentation.")
    parser.add_argument("input", help="Path to input video file")
    parser.add_argument("output", help="Path to output PPTX file")
    parser.add_argument("--threshold", type=int, default=2000, help="Sensitivity threshold (lower=more slides)")
    parser.add_argument("--rate", type=float, default=1.0, help="Sample rate in seconds")
    
    args = parser.parse_args()
    convert_mp4_to_pptx(args.input, args.output, args.threshold, args.rate)

if __name__ == "__main__":
    main()
